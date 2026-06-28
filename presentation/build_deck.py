from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "presentation"
PREVIEW_DIR = OUT_DIR / "previews"
PPTX_PATH = OUT_DIR / "HSL_Technical_Assessment_3DGS_Composition.pptx"
FINAL_RENDER_PATH = ROOT / "hsl_final_result.png"

W, H = 1920, 1080
SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5

COLORS = {
    "paper": "F7F4EC",
    "ink": "111827",
    "muted": "5F6B73",
    "teal": "1F766F",
    "blue": "365E9D",
    "coral": "D7613D",
    "gold": "D9A441",
    "line": "C8C2B4",
    "soft": "EDE7DA",
    "white": "FFFFFF",
}


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def px_to_in(x: float, y: float, w: float, h: float) -> tuple[float, float, float, float]:
    return (
        x / W * SLIDE_W_IN,
        y / H * SLIDE_H_IN,
        w / W * SLIDE_W_IN,
        h / H * SLIDE_H_IN,
    )


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int,
    color: str = "ink",
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = "Aptos",
    line_spacing: float | None = None,
):
    left, top, width, height = [Inches(v) for v in px_to_in(x, y, w, h)]
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for index, line in enumerate(lines):
        para = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        para.text = line
        para.alignment = align
        if line_spacing is not None:
            para.line_spacing = line_spacing
        run = para.runs[0]
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(COLORS[color])
    return box


def add_rect(slide, x, y, w, h, fill, line=None):
    left, top, width, height = [Inches(v) for v in px_to_in(x, y, w, h)]
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[fill])
    if line:
        shape.line.color.rgb = rgb(COLORS[line])
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_circle(slide, x, y, d, fill, transparency=0):
    left, top, width, height = [Inches(v) for v in px_to_in(x, y, d, d)]
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[fill])
    shape.fill.transparency = transparency
    shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color="line", width=2):
    left, top, w, h = [Inches(v) for v in px_to_in(x1, y1, x2 - x1, y2 - y1)]
    line = slide.shapes.add_connector(1, left, top, left + w, top + h)
    line.line.color.rgb = rgb(COLORS[color])
    line.line.width = Pt(width)
    return line


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float):
    if not path.exists():
        return None
    image = Image.open(path)
    img_w, img_h = image.size
    frame_ratio = w / h
    img_ratio = img_w / img_h
    if img_ratio >= frame_ratio:
        draw_w = w
        draw_h = w / img_ratio
    else:
        draw_h = h
        draw_w = h * img_ratio
    draw_x = x + (w - draw_w) / 2
    draw_y = y + (h - draw_h) / 2
    left, top, width, height = [Inches(v) for v in px_to_in(draw_x, draw_y, draw_w, draw_h)]
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def background(slide):
    add_rect(slide, 0, 0, W, H, "paper")


def pptx_splat_cloud(slide, origin_x, origin_y, scale=1.0):
    scene_points = [
        (0, 30, 20),
        (70, 0, 16),
        (145, 24, 22),
        (210, 3, 13),
        (285, 42, 18),
        (345, 16, 26),
        (420, 36, 15),
        (38, 118, 26),
        (115, 92, 18),
        (185, 135, 22),
        (255, 105, 14),
        (330, 130, 21),
        (400, 98, 17),
        (65, 215, 18),
        (140, 250, 26),
        (235, 220, 16),
        (310, 260, 23),
        (395, 230, 18),
        (100, 345, 22),
        (210, 360, 16),
        (325, 340, 24),
    ]
    asset_points = [
        (230, 155, 32),
        (275, 145, 22),
        (315, 175, 28),
        (250, 205, 24),
        (295, 225, 34),
        (345, 210, 20),
        (268, 272, 18),
        (325, 285, 23),
    ]
    for x, y, d in scene_points:
        add_circle(slide, origin_x + x * scale, origin_y + y * scale, d * scale, "blue", 24)
    for x, y, d in asset_points:
        add_circle(slide, origin_x + x * scale, origin_y + y * scale, d * scale, "coral", 8)


def slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    add_rect(slide, 0, 0, 34, H, "teal")
    add_text(slide, "Human Sensing Lab\nTechnical Assessment", 118, 92, 520, 80, 22, "muted")
    add_text(
        slide,
        "Composing\n3DGS Assets",
        116,
        190,
        720,
        230,
        66,
        "ink",
        bold=True,
        font="Aptos Display",
        line_spacing=0.85,
    )
    add_text(
        slide,
        "Separately train Bicycle and a synthetic chair, then merge the foreground Gaussians into the scene.",
        122,
        470,
        700,
        110,
        24,
        "muted",
    )
    add_text(slide, "Bicycle scene + synthetic chair", 124, 740, 520, 42, 20, "teal", bold=True)
    add_text(slide, "June 2026", 124, 790, 240, 36, 18, "muted")
    add_picture_contain(slide, FINAL_RENDER_PATH, 930, 160, 790, 525)
    add_rect(slide, 930, 160, 790, 525, "white", "line").fill.transparency = 100
    add_text(slide, "final composed render", 930, 710, 390, 34, 19, "teal", bold=True)
    add_text(slide, "one merged point_cloud.ply", 930, 790, 620, 44, 27, "ink", bold=True)
    add_text(slide, "Rendered with Bicycle cameras for shared depth and alpha ordering.", 930, 840, 700, 70, 21, "muted")


def slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    add_text(slide, "Composition Method", 110, 82, 900, 78, 48, "ink", bold=True, font="Aptos Display")
    add_text(
        slide,
        "The foreground asset is transformed in 3DGS parameter space, then merged before rendering.",
        112,
        150,
        1180,
        56,
        22,
        "muted",
    )
    add_text(slide, "p_scene = s R p_asset + t", 250, 280, 1040, 80, 45, "ink", bold=True, align=PP_ALIGN.CENTER)
    add_line(slide, 300, 395, 1520, 395, "line", 3)
    steps = [
        (250, "blue", "Train", "Bicycle scene and chair asset as separate 3DGS models."),
        (790, "gold", "Transform", "Apply scale, rotation, and translation to every foreground Gaussian."),
        (1330, "coral", "Merge", "Write one composed point_cloud.ply and render with Bicycle cameras."),
    ]
    for x, color, label, body in steps:
        add_circle(slide, x, 470, 76, color, 0)
        add_text(slide, label, x - 145, 575, 370, 42, 29, "ink", bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x - 190, 635, 470, 100, 22, "muted", align=PP_ALIGN.CENTER)
    add_text(slide, "Gaussian updates", 160, 820, 470, 40, 24, "teal", bold=True)
    add_text(
        slide,
        "center: similarity transform\nscale: add log(s)\nrotation: q_place times q_asset",
        160,
        875,
        650,
        110,
        21,
        "ink",
    )
    add_text(slide, "Why this helps", 1030, 820, 430, 40, 24, "coral", bold=True)
    add_text(
        slide,
        "Scene and asset splats are depth-sorted together, avoiding a separate 2D paste step.",
        1030,
        875,
        690,
        95,
        21,
        "ink",
    )


def slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    add_text(slide, "Result & Defense", 110, 82, 850, 78, 48, "ink", bold=True, font="Aptos Display")
    add_text(
        slide,
        "The output is not a 2D paste. It is one merged 3DGS point cloud rendered through the Bicycle cameras.",
        112,
        150,
        1260,
        56,
        22,
        "muted",
    )
    add_picture_contain(slide, FINAL_RENDER_PATH, 110, 260, 900, 595)
    add_rect(slide, 110, 260, 900, 595, "white", "line").fill.transparency = 100
    add_text(slide, "final composed render", 110, 878, 370, 38, 19, "teal", bold=True)
    rows = [
        (
            280,
            "Result",
            "Separate Bicycle and Chair 3DGS models\nmerged into one splat point cloud.",
            "teal",
        ),
        (
            455,
            "Scale reasoning",
            "Cross-dataset scale is ambiguous;\nbounds and visual anchors define s, R, t.",
            "blue",
        ),
        (
            630,
            "Conflicts",
            "One rasterizer pass gives depth sorting;\nfloaters and soft splats remain visible.",
            "coral",
        ),
        (
            805,
            "Extension idea",
            "Optimize asset-only SH/exposure\nwhile keeping geometry fixed.",
            "gold",
        ),
    ]
    for y, title, body, color in rows:
        add_circle(slide, 1072, y - 4, 46, color, 0)
        add_text(slide, title, 1144, y - 14, 500, 42, 28, "ink", bold=True)
        add_text(slide, body, 1144, y + 40, 620, 104, 21, "muted")
    add_text(
        slide,
        "Submission assets: code, process notes, Kaggle/Colab runners, final render, and this 3-slide deck.",
        110,
        960,
        1460,
        34,
        18,
        "muted",
    )


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    prs.save(PPTX_PATH)


def load_font(size: int, bold: bool = False):
    candidates = [
        "C:/Windows/Fonts/aptos.ttf",
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
    ]
    for candidate in candidates:
        path = Path(candidate)
        if path.exists():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()


def draw_text(draw, xy, text, size, color="ink", bold=False, anchor=None, spacing=8):
    font = load_font(size, bold)
    draw.multiline_text(
        xy,
        text,
        fill="#" + COLORS[color],
        font=font,
        spacing=spacing,
        anchor=anchor,
        align="center" if anchor == "mm" else "left",
    )


def draw_cloud(draw, origin_x, origin_y, scale=1.0):
    scene_points = [
        (0, 30, 20), (70, 0, 16), (145, 24, 22), (210, 3, 13), (285, 42, 18),
        (345, 16, 26), (420, 36, 15), (38, 118, 26), (115, 92, 18),
        (185, 135, 22), (255, 105, 14), (330, 130, 21), (400, 98, 17),
        (65, 215, 18), (140, 250, 26), (235, 220, 16), (310, 260, 23),
        (395, 230, 18), (100, 345, 22), (210, 360, 16), (325, 340, 24),
    ]
    asset_points = [
        (230, 155, 32), (275, 145, 22), (315, 175, 28), (250, 205, 24),
        (295, 225, 34), (345, 210, 20), (268, 272, 18), (325, 285, 23),
    ]
    for x, y, d in scene_points:
        xx, yy, dd = origin_x + x * scale, origin_y + y * scale, d * scale
        draw.ellipse([xx, yy, xx + dd, yy + dd], fill="#365E9D")
    for x, y, d in asset_points:
        xx, yy, dd = origin_x + x * scale, origin_y + y * scale, d * scale
        draw.ellipse([xx, yy, xx + dd, yy + dd], fill="#D7613D")


def paste_image_contain(canvas: Image.Image, path: Path, box: tuple[int, int, int, int]) -> None:
    x, y, w, h = box
    draw = ImageDraw.Draw(canvas)
    draw.rectangle([x, y, x + w, y + h], fill="#" + COLORS["white"], outline="#" + COLORS["line"], width=2)
    if not path.exists():
        draw_text(draw, (x + 40, y + h // 2), f"{path.name} missing", 34, "muted")
        return
    image = Image.open(path).convert("RGB")
    img_ratio = image.width / image.height
    frame_ratio = w / h
    if img_ratio >= frame_ratio:
        new_w = w
        new_h = round(new_w / img_ratio)
    else:
        new_h = h
        new_w = round(new_h * img_ratio)
    image = image.resize((new_w, new_h), Image.Resampling.LANCZOS)
    canvas.paste(image, (x + (w - new_w) // 2, y + (h - new_h) // 2))


def preview_slide_1(path: Path) -> None:
    img = Image.new("RGB", (W, H), "#" + COLORS["paper"])
    draw = ImageDraw.Draw(img)
    draw.rectangle([0, 0, 34, H], fill="#" + COLORS["teal"])
    draw_text(draw, (118, 92), "Human Sensing Lab\nTechnical Assessment", 32, "muted")
    draw_text(draw, (116, 190), "Composing\n3DGS Assets", 84, "ink", True, spacing=0)
    draw_text(draw, (122, 470), "Separately train Bicycle and a synthetic chair,\nthen merge the foreground Gaussians into the scene.", 30, "muted")
    draw_text(draw, (124, 740), "Bicycle scene + synthetic chair", 28, "teal", True)
    draw_text(draw, (124, 790), "June 2026", 25, "muted")
    paste_image_contain(img, FINAL_RENDER_PATH, (930, 160, 790, 525))
    draw_text(draw, (930, 710), "final composed render", 27, "teal", True)
    draw_text(draw, (930, 790), "one merged point_cloud.ply", 36, "ink", True)
    draw_text(draw, (930, 840), "Rendered with Bicycle cameras for shared\ndepth and alpha ordering.", 28, "muted")
    img.save(path)


def preview_slide_2(path: Path) -> None:
    img = Image.new("RGB", (W, H), "#" + COLORS["paper"])
    draw = ImageDraw.Draw(img)
    draw_text(draw, (110, 82), "Composition Method", 62, "ink", True)
    draw_text(draw, (112, 150), "The foreground asset is transformed in 3DGS parameter space, then merged before rendering.", 30, "muted")
    draw_text(draw, (960, 320), "p_scene = s R p_asset + t", 58, "ink", True, anchor="mm")
    draw.line([300, 395, 1520, 395], fill="#" + COLORS["line"], width=5)
    steps = [
        (250, "blue", "Train", "Bicycle scene and chair asset\nas separate 3DGS models."),
        (790, "gold", "Transform", "Apply scale, rotation, and translation\nto every foreground Gaussian."),
        (1330, "coral", "Merge", "Write one composed point_cloud.ply\nand render with Bicycle cameras."),
    ]
    for x, color, label, body in steps:
        draw.ellipse([x, 470, x + 76, 546], fill="#" + COLORS[color])
        draw_text(draw, (x + 38, 508), str(steps.index((x, color, label, body)) + 1), 29, "white", True, anchor="mm")
        draw_text(draw, (x - 145, 575), label, 38, "ink", True)
        draw_text(draw, (x - 190, 635), body, 28, "muted")
    draw_text(draw, (160, 820), "Gaussian updates", 32, "teal", True)
    draw_text(draw, (160, 875), "center: similarity transform\nscale: add log(s)\nrotation: q_place times q_asset", 28, "ink")
    draw_text(draw, (1030, 820), "Why this helps", 32, "coral", True)
    draw_text(draw, (1030, 875), "Scene and asset splats are depth-sorted together,\navoiding a separate 2D paste step.", 28, "ink")
    img.save(path)


def preview_slide_3(path: Path) -> None:
    img = Image.new("RGB", (W, H), "#" + COLORS["paper"])
    draw = ImageDraw.Draw(img)
    draw_text(draw, (110, 82), "Result & Defense", 62, "ink", True)
    draw_text(draw, (112, 150), "The output is not a 2D paste. It is one merged 3DGS point cloud rendered through the Bicycle cameras.", 30, "muted")
    paste_image_contain(img, FINAL_RENDER_PATH, (110, 260, 900, 595))
    draw_text(draw, (110, 878), "final composed render", 27, "teal", True)
    rows = [
        (280, "Result", "Separate Bicycle and Chair 3DGS models\nmerged into one splat point cloud.", "teal"),
        (455, "Scale reasoning", "Cross-dataset scale is ambiguous;\nbounds and visual anchors define s, R, t.", "blue"),
        (630, "Conflicts", "One rasterizer pass gives depth sorting;\nfloaters and soft splats remain visible.", "coral"),
        (805, "Extension idea", "Optimize asset-only SH/exposure\nwhile keeping geometry fixed.", "gold"),
    ]
    for y, title, body, color in rows:
        draw.ellipse([1072, y - 4, 1118, y + 42], fill="#" + COLORS[color])
        draw_text(draw, (1144, y - 14), title, 38, "ink", True)
        draw_text(draw, (1144, y + 38), body, 29, "muted")
    draw_text(draw, (110, 960), "Submission assets: code, process notes, Kaggle/Colab runners, final render, and this 3-slide deck.", 25, "muted")
    img.save(path)


def build_previews() -> None:
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    preview_fns = [preview_slide_1, preview_slide_2, preview_slide_3]
    paths = []
    for index, fn in enumerate(preview_fns, start=1):
        path = PREVIEW_DIR / f"slide_{index}.png"
        fn(path)
        paths.append(path)

    montage = Image.new("RGB", (W, H * 3), "#" + COLORS["paper"])
    for index, path in enumerate(paths):
        montage.paste(Image.open(path), (0, H * index))
    montage.save(PREVIEW_DIR / "montage.png")


def inspect_pptx_package() -> None:
    with ZipFile(PPTX_PATH) as zf:
        names = zf.namelist()
        slide_xml = [name for name in names if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
        if len(slide_xml) != 3:
            raise RuntimeError(f"Expected 3 slides, found {len(slide_xml)}")
        for name in slide_xml:
            xml = zf.read(name).decode("utf-8", errors="ignore")
            if "Slide Number" in xml or "sldNum" in xml:
                raise RuntimeError(f"Unexpected slide-number placeholder in {name}")


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    build_pptx()
    build_previews()
    inspect_pptx_package()
    print(f"Wrote {PPTX_PATH}")
    print(f"Wrote previews to {PREVIEW_DIR}")


if __name__ == "__main__":
    main()
